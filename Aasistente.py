from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import google.generativeai as genai
import io
import sqlite3
import os
from pptx import Presentation
from pypdf import PdfReader

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], 
    allow_credentials=True,
    allow_methods=["*"], 
    allow_headers=["*"], 
)

# Llave de Gemini desde Variables de Entorno
API_KEY = os.environ.get("GEMINI_API_KEY", "TU_LLAVE_TEMPORAL_AQUI") 
genai.configure(api_key=API_KEY)
model = genai.GenerativeModel('gemini-2.5-flash')

# --- MODELOS DE DATOS ---
class UserAuth(BaseModel):
    username: str
    password: str

class ClaseData(BaseModel):
    usuario_id: int
    nombre: str
    dia: str
    inicio: str
    fin: str

class EstudiarTemaRequest(BaseModel):
    usuario_id: int
    curso: str
    tema: str
    tipo_output: str

# --- BASE DE DATOS LOCAL CON SOPORTE MULTIUSUARIO ---
def iniciar_memoria():
    conexion = sqlite3.connect('nami_memoria.db')
    cursor = conexion.cursor()
    
    # 1. Tabla de Usuarios (NUEVA)
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS usuarios (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT UNIQUE,
            password TEXT
        )
    ''')
    
    # 2. Tabla de Horario (Ahora requiere usuario_id)
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS clases (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            usuario_id INTEGER,
            nombre TEXT,
            dia TEXT,
            inicio TEXT,
            fin TEXT
        )
    ''')
    
    # 3. Tabla de Biblioteca (Ahora requiere usuario_id)
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS biblioteca (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            usuario_id INTEGER,
            curso TEXT,
            tema TEXT,
            contenido TEXT
        )
    ''')
    conexion.commit()
    conexion.close()

iniciar_memoria()

# --- FUNCIONES AUXILIARES ---
def extraer_texto_pptx(file_content):
    prs = Presentation(io.BytesIO(file_content))
    return " ".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])

def extraer_texto_pdf(file_content):
    reader = PdfReader(io.BytesIO(file_content))
    return " ".join([page.extract_text() for page in reader.pages])


# --- ENDPOINTS DE AUTENTICACIÓN (LOGIN / REGISTRO) ---

@app.post("/registro")
async def registrar_usuario(user: UserAuth):
    conexion = sqlite3.connect('nami_memoria.db')
    cursor = conexion.cursor()
    try:
        cursor.execute('INSERT INTO usuarios (username, password) VALUES (?, ?)', (user.username.lower(), user.password))
        usuario_id = cursor.lastrowid
        conexion.commit()
        return {"mensaje": "Usuario creado con éxito", "usuario_id": usuario_id, "username": user.username}
    except sqlite3.IntegrityError:
        raise HTTPException(status_code=400, detail="El nombre de usuario ya existe. Elige otro.")
    finally:
        conexion.close()

@app.post("/login")
async def iniciar_sesion(user: UserAuth):
    conexion = sqlite3.connect('nami_memoria.db')
    cursor = conexion.cursor()
    cursor.execute('SELECT id, username FROM usuarios WHERE username = ? AND password = ?', (user.username.lower(), user.password))
    usuario = cursor.fetchone()
    conexion.close()
    
    if usuario:
        return {"mensaje": "Login exitoso", "usuario_id": usuario[0], "username": usuario[1]}
    else:
        raise HTTPException(status_code=401, detail="Usuario o contraseña incorrectos")


# --- ENDPOINTS DE HORARIO (AHORA SON PRIVADOS) ---

@app.post("/agregar-clase")
async def agregar_clase(clase: ClaseData):
    conexion = sqlite3.connect('nami_memoria.db')
    cursor = conexion.cursor()
    cursor.execute('INSERT INTO clases (usuario_id, nombre, dia, inicio, fin) VALUES (?, ?, ?, ?, ?)', 
                   (clase.usuario_id, clase.nombre, clase.dia, clase.inicio, clase.fin))
    clase_id = cursor.lastrowid
    conexion.commit()
    conexion.close()
    return {"id": clase_id}

@app.get("/cargar-horario/{usuario_id}")
async def cargar_horario(usuario_id: int):
    conexion = sqlite3.connect('nami_memoria.db')
    cursor = conexion.cursor()
    cursor.execute('SELECT id, nombre, dia, inicio, fin FROM clases WHERE usuario_id = ?', (usuario_id,))
    clases = [{"id": row[0], "nombre": row[1], "dia": row[2], "inicio": row[3], "fin": row[4]} for row in cursor.fetchall()]
    conexion.close()
    return clases

@app.delete("/eliminar-clase/{clase_id}")
async def eliminar_clase(clase_id: int):
    conexion = sqlite3.connect('nami_memoria.db')
    cursor = conexion.cursor()
    cursor.execute('DELETE FROM clases WHERE id = ?', (clase_id,))
    conexion.commit()
    conexion.close()
    return {"mensaje": "Eliminada"}


# --- ENDPOINT IA: MODO ESTUDIO LIBRE (Público, no guarda datos) ---

@app.post("/procesar-material")
async def procesar_material(file: UploadFile = File(...), tipo_output: str = "quiz"):
    try:
        content = await file.read()
        texto_extraido = extraer_texto_pptx(content) if file.filename.endswith(".pptx") else extraer_texto_pdf(content) if file.filename.endswith(".pdf") else ""
        if not texto_extraido: raise HTTPException(status_code=400, detail="Error de formato")
        instrucciones = {"quiz": "Crea un cuestionario de 5 preguntas de opción múltiple con respuestas...", "resumen": "Haz un resumen estructurado..."}
        prompt = f"{instrucciones.get(tipo_output, 'Resume esto')}: {texto_extraido[:10000]}"
        response = model.generate_content(prompt)
        return {"resultado": response.text}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# --- ENDPOINTS DE BIBLIOTECA (AHORA SON PRIVADOS) ---

@app.post("/subir-material-tema")
async def subir_material_tema(usuario_id: int = Form(...), curso: str = Form(...), tema: str = Form(...), file: UploadFile = File(...)):
    try:
        content = await file.read()
        texto_extraido = extraer_texto_pptx(content) if file.filename.endswith(".pptx") else extraer_texto_pdf(content) if file.filename.endswith(".pdf") else ""
        if not texto_extraido: raise HTTPException(status_code=400, detail="El archivo está vacío o no es soportado.")

        conexion = sqlite3.connect('nami_memoria.db')
        cursor = conexion.cursor()
        cursor.execute('INSERT INTO biblioteca (usuario_id, curso, tema, contenido) VALUES (?, ?, ?, ?)', (usuario_id, curso, tema, texto_extraido))
        conexion.commit()
        conexion.close()
        return {"mensaje": f"Material guardado con éxito en {curso} - {tema}"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/obtener-temas/{usuario_id}/{curso}")
async def obtener_temas(usuario_id: int, curso: str):
    conexion = sqlite3.connect('nami_memoria.db')
    cursor = conexion.cursor()
    cursor.execute('SELECT DISTINCT tema FROM biblioteca WHERE usuario_id = ? AND curso = ?', (usuario_id, curso))
    temas = [row[0] for row in cursor.fetchall()]
    conexion.close()
    return {"temas": temas}

@app.post("/estudiar-tema")
async def estudiar_tema(req: EstudiarTemaRequest):
    conexion = sqlite3.connect('nami_memoria.db')
    cursor = conexion.cursor()
    cursor.execute('SELECT contenido FROM biblioteca WHERE usuario_id = ? AND curso = ? AND tema = ?', (req.usuario_id, req.curso, req.tema))
    filas = cursor.fetchall()
    conexion.close()

    if not filas:
        raise HTTPException(status_code=404, detail="No hay material guardado en este tema.")

    texto_completo = " ".join([fila[0] for fila in filas])

    instrucciones = {
        "quiz": "Crea un cuestionario nivel universitario de 5 preguntas de opción múltiple (con respuestas al final) basado estrictamente en el siguiente texto.",
        "resumen": "Haz un resumen estructurado, con puntos clave y viñetas, basado estrictamente en el siguiente texto."
    }
    
    prompt = f"{instrucciones.get(req.tipo_output, 'Resume esto')}: {texto_completo[:15000]}"
    
    try:
        response = model.generate_content(prompt)
        return {"resultado": response.text}
    except Exception as e:
        raise HTTPException(status_code=500, detail="Error conectando con NAMI: " + str(e))
